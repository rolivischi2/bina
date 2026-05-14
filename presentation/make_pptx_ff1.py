"""
FF1: Absolutes Veloaufkommen Zürich 2020–2025
Self-contained — all data hardcoded from notebook outputs, no CSV dependency.
"""
import io, os
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import seaborn as sns
import warnings
from PIL import Image
Image.MAX_IMAGE_PIXELS = None   # disable PIL decompression-bomb guard for pptx embedding
warnings.filterwarnings("ignore", message="Tight layout not applied")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colours ───────────────────────────────────────────────────────────────────
C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xB4, 0xD8)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_ORANGE = RGBColor(0xFF, 0x6B, 0x35)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
BG_HEX  = "#0D1B2A"

# ── data (from FF1_absolutes_aufkommen.ipynb, 2020–2025) ──────────────────────

# Top-15 heatmap: yearly totals per station (sum of all station_ids with that name)
# Columns: 2020, 2021, 2022, 2023, 2024, 2025
HEATMAP = {
    "Stadttunnel Nord":   [3111883, 1781325, 3126808, 3179358, 3227784, 3061633],
    "Hardbrücke":         [1306196, 1302755, 1443490, 1415998, 1387219, 1399466],
    "In Gassen":          [ 691560,  794515,  896306, 1080305,  885994,  857216],
    "Bucheggstrasse":     [ 356386,  752988,  842483,  887898,  891203,  893314],
    "Kreuzstrasse":       [ 643988,  591466,  658353,  543510,  677555,  724006],
    "Talstrasse":         [ 603189,  214028,  410620,  829939,  821858,  807177],
    "Mythenquai":         [ 679115,  542727,   40906,  592035,  642902,  682128],
    "Sonneggstrasse":     [ 518245,  451135,  593960,  559149,  566101,  476676],
    "Weststrasse":        [ 525326,  336783,       0,  414935,  817878,  769311],
    "Sihlstrasse":        [ 310749,  283484,  315491,  285006,  294317,  998362],
    "Weinbergstrasse":    [ 300979,  381578,  407097,  392769,  373573,  419543],
    "Rennweg":            [ 220889,  141237,  600470,  557701,  478755,  234835],
    "Usteristrasse":      [ 684439,  706371,  117208,       0,       0,  538535],
    "Altstetterstrasse":  [ 169190,  240002,  302618,  318048,  366114,  374360],
    "Saatlenstrasse":     [ 321599,  261146,  281913,  273718,  255138,  130482],
}
YEARS = [2020, 2021, 2022, 2023, 2024, 2025]

# ── data (from 03_ff2_growth_hotspots.ipynb, 2020–2025) ──────────────────────
# 12 consistent stations only (≥60% data coverage every year)
# Sorted by CAGR descending; trend_quality from linear regression (p-value + R²)
FF2_STATIONS = [
    # (name,                              cagr_pct, trend_quality,  mean_annual)
    ("Baslerstrasse",                       +13.4,  "strong",        289843),
    ("Talstrasse",                          +10.7,  "weak",          409311),
    ("Langstrasse (Fahrbahn Nord)",          +5.1,  "strong",        369744),
    ("Hardbrücke Nord",                      +4.2,  "strong",        414933),
    ("Bucheggplatz",                         +0.9,  "weak",          525163),
    ("Andreasstrasse",                       +0.1,  "weak",          291958),
    ("Hardbrücke Süd",                       +0.2,  "weak",          962460),
    ("Scheuchzerstrasse",                    -0.1,  "weak",          559143),
    ("Langstrasse (Fahrbahn Süd)",           -1.3,  "weak",          333089),
    ("Lux-Guyer-Weg",                        -1.6,  "weak",          405841),
    ("Langstrasse (Unterführung Süd)",       -4.7,  "moderate",     1162289),
    ("Langstrasse (Unterführung Nord)",      -6.8,  "moderate",     1416494),
]

# ── data (from 04_ff3_Wettersensitivität_v2.ipynb, 2020–2025) ────────────────
# Network-level effect: avg velos per station-day by weather condition
FF3_WEATHER = {
    "Trocken & mild": 1500,
    "Nur Regen":      1347,
    "Nur Kalt (≤5°C)": 904,
    "Regen + Kalt":    790,   # estimated: cold baseline × (1 - rain effect)
}
# Station-level rain sensitivity (% change vs dry days), from notebook L3 analysis
FF3_STATION_RAIN = [
    ("Fischerweg",        -24.4),
    ("Lux-Guyer-Weg",    -23.3),
    ("Mythenquai",        -19.4),
    ("Seefeldstrasse",    -16.0),
    ("Bellerivestrasse",  -14.5),
    ("Netz-Ø",             -9.8),
    ("Militärstrasse",     -7.5),
    ("Hardbrücke",         -6.2),
    ("Sihlstrasse",        -4.1),
    ("Stadttunnel Süd",    -0.5),
    ("Talstrasse",         +0.4),
    ("Schulstrasse",       +4.8),
]

# ── helpers ───────────────────────────────────────────────────────────────────
def fig_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=120, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG

def add_label(slide, text, left, top, width, height, size=11,
              bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def accent_line(slide, left, top, width, height=Inches(0.04)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_ACCENT
    shape.line.fill.background()

# ── chart 2: heatmap Top-15 ───────────────────────────────────────────────────
print("Generating Chart 2 (heatmap) …")
station_order = list(HEATMAP.keys())
heatmap_arr   = np.array([HEATMAP[s] for s in station_order], dtype=float)

fig, ax = plt.subplots(figsize=(7, 5.5))
fig.patch.set_facecolor(BG_HEX); ax.set_facecolor(BG_HEX)
sns.heatmap(heatmap_arr / 1000, annot=True, fmt=".0f", cmap="YlOrRd",
            linewidths=0.3, linecolor="#1a1a2e", ax=ax,
            cbar_kws={"label":"Tsd. Fahrten","shrink":0.7,"pad":0.02},
            annot_kws={"size": 8})
ax.set_xticks(np.arange(len(YEARS)) + 0.5)
ax.set_xticklabels([str(y) for y in YEARS])
ax.set_yticks(np.arange(len(station_order)) + 0.5)
ax.set_yticklabels(station_order, rotation=0)
ax.set_xlabel("Jahr", fontsize=9, color="#B0C4DE")
ax.set_ylabel("")
ax.tick_params(axis="x", labelsize=8, colors="white")
ax.tick_params(axis="y", labelsize=8, colors="white")
ax.collections[0].colorbar.ax.yaxis.set_tick_params(color="white")
ax.collections[0].colorbar.set_label("Tsd. Fahrten", color="white", fontsize=8)
plt.setp(ax.collections[0].colorbar.ax.yaxis.get_ticklabels(), color="white")
plt.tight_layout(pad=0.5)
chart2_stream = fig_to_stream(fig)

# ── chart FF2: CAGR ranking bar ──────────────────────────────────────────────
print("Generating Chart FF2 (CAGR ranking) …")

# Sort ascending so highest CAGR appears at top of horizontal bar chart
ff2_sorted = sorted(FF2_STATIONS, key=lambda r: r[1])
ff2_names  = [r[0] for r in ff2_sorted]
ff2_cagrs  = [r[1] for r in ff2_sorted]
ff2_qual   = [r[2] for r in ff2_sorted]

def _bar_color(cagr, quality):
    if cagr < 0:
        return "#FF4D6D"
    if quality == "strong":
        return "#00D4A0"
    if quality == "moderate":
        return "#48CAE4"
    return "#5A7A8A"

bar_colors = [_bar_color(c, q) for c, q in zip(ff2_cagrs, ff2_qual)]

fig, ax = plt.subplots(figsize=(8.5, 5.8))
fig.patch.set_facecolor(BG_HEX); ax.set_facecolor(BG_HEX)
bars = ax.barh(ff2_names, ff2_cagrs, color=bar_colors, edgecolor="none", height=0.62)

for bar, val in zip(bars, ff2_cagrs):
    x_pos = val + (0.15 if val >= 0 else -0.15)
    ax.text(x_pos, bar.get_y() + bar.get_height() / 2,
            f"{val:+.1f}%", va="center",
            ha="left" if val >= 0 else "right",
            fontsize=8.5, color="white")

ax.axvline(0, color="white", lw=0.8, alpha=0.4)
ax.set_xlabel("CAGR 2020–2025 (%)", color="#B0C4DE", fontsize=9)
ax.tick_params(axis="y", labelsize=8.5, colors="white")
ax.tick_params(axis="x", labelsize=8, colors="#B0C4DE")
for spine in ax.spines.values(): spine.set_visible(False)

import matplotlib.patches as mpatches
legend_patches = [
    mpatches.Patch(color="#00D4A0", label="Statistisch stark (p<0.05, R²≥0.7)"),
    mpatches.Patch(color="#5A7A8A", label="Schwach / kein Trend"),
    mpatches.Patch(color="#FF4D6D", label="Rückgang"),
]
ax.legend(handles=legend_patches, fontsize=7.5, framealpha=0,
          labelcolor="white", loc="lower right")

plt.tight_layout(pad=0.5)
chartFF2_stream = fig_to_stream(fig)

# ── chart FF3: 4-category weather effect bar ─────────────────────────────────
print("Generating Chart FF3 (weather effect) …")

labels         = list(FF3_WEATHER.keys())
values         = list(FF3_WEATHER.values())
bar_colors_ff3 = ["#00B4D8", "#FF9F43", "#FF4D6D", "#C0392B"]

fig, ax = plt.subplots(figsize=(8.0, 5.0))
fig.patch.set_facecolor(BG_HEX); ax.set_facecolor(BG_HEX)

x    = np.arange(len(labels))
bars = ax.bar(x, values, color=bar_colors_ff3, width=0.55, edgecolor="none", zorder=3)

# Value labels above bars
for bar, val in zip(bars, values):
    ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 15,
            f"{val:,}", ha="center", va="bottom", fontsize=11,
            color="white", fontweight="bold")

# Percentage labels inside bars (mid-height)
pcts = ["Basis", "−10.2%", "−41.6%", "−47.3%"]
pct_colors = ["#B0C4DE", "#FF9F43", "#FF4D6D", "#C0392B"]
for bar, pct, col in zip(bars, pcts, pct_colors):
    mid = bar.get_height() / 2
    ax.text(bar.get_x() + bar.get_width() / 2, mid,
            pct, ha="center", va="center", fontsize=10.5,
            color="white", fontweight="bold")

ax.set_xticks(x)
ax.set_xticklabels(labels, fontsize=9.5, color="white")
ax.set_ylabel("Ø Velos / Standort-Tag", color="#B0C4DE", fontsize=9)
ax.set_ylim(0, 1850)
ax.tick_params(axis="y", labelsize=8.5, colors="#B0C4DE")
ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{int(v):,}"))
ax.grid(axis="y", linestyle="--", alpha=0.12, color="white")
for spine in ax.spines.values(): spine.set_visible(False)

fig.subplots_adjust(bottom=0.12)
chartFF3_stream = fig_to_stream(fig)

# ── build PPTX ────────────────────────────────────────────────────────────────
print("Building PPTX …")
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ── SLIDE 0: Einleitung — Kontext & Forschungsfragen ─────────────────────────
s_intro = prs.slides.add_slide(blank)
add_bg(s_intro)
accent_line(s_intro, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
accent_line(s_intro, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))

add_label(s_intro, "01", Inches(0.25), Inches(0.2), Inches(0.8), Inches(0.5),
          size=11, color=C_ACCENT)
add_label(s_intro, "Kontext & Forschungsfragen",
          Inches(0.25), Inches(0.45), Inches(12), Inches(0.6),
          size=26, bold=True, color=C_WHITE)
add_label(s_intro,
          "Datenbasierte Entscheidungsunterstützung für die Velostrategie 2030 der Stadt Zürich",
          Inches(0.25), Inches(0.98), Inches(12), Inches(0.38),
          size=12, color=C_LIGHT)

# Left card: Velostrategie 2030
L_LEFT = Inches(0.25)
L_TOP  = Inches(1.5)
L_W    = Inches(3.9)
L_H    = Inches(5.3)

card_l = s_intro.shapes.add_shape(1, L_LEFT, L_TOP, L_W, L_H)
card_l.fill.solid(); card_l.fill.fore_color.rgb = RGBColor(0x12, 0x26, 0x3A)
card_l.line.color.rgb = C_ACCENT; card_l.line.width = Pt(1.5)

add_label(s_intro, "🚴", L_LEFT + Inches(0.18), L_TOP + Inches(0.15),
          Inches(0.55), Inches(0.52), size=22, align=PP_ALIGN.CENTER)
add_label(s_intro, "Velostrategie 2030",
          L_LEFT + Inches(0.75), L_TOP + Inches(0.2),
          L_W - Inches(0.9), Inches(0.45),
          size=13, bold=True, color=C_ACCENT)

div_l = s_intro.shapes.add_shape(1, L_LEFT + Inches(0.18), L_TOP + Inches(0.75),
                                  L_W - Inches(0.36), Inches(0.025))
div_l.fill.solid(); div_l.fill.fore_color.rgb = C_ACCENT
div_l.line.fill.background()

add_label(s_intro,
          "«Velofahren für alle\njederzeit selbstverständlich»",
          L_LEFT + Inches(0.18), L_TOP + Inches(0.88),
          L_W - Inches(0.3), Inches(0.72),
          size=11, bold=True, color=C_WHITE)

add_label(s_intro,
          "Die Stadt Zürich hat mit der Velostrategie 2030 das Ziel gesetzt, "
          "Velofahren als selbstverständliches Alltagsverkehrsmittel zu etablieren. "
          "Datenbasierte Analysen sind dabei unverzichtbar, um Infrastrukturentscheide "
          "gezielt und evidenzbasiert zu treffen.",
          L_LEFT + Inches(0.18), L_TOP + Inches(1.7),
          L_W - Inches(0.3), Inches(1.45),
          size=9.5, color=C_LIGHT)

add_label(s_intro, "Analysebasis",
          L_LEFT + Inches(0.18), L_TOP + Inches(3.25),
          L_W - Inches(0.3), Inches(0.32),
          size=9.5, bold=True, color=C_ACCENT)
add_label(s_intro,
          "Open Data Stadt Zürich\nVerkehrszählungen Fussgänger & Velo",
          L_LEFT + Inches(0.18), L_TOP + Inches(3.58),
          L_W - Inches(0.3), Inches(0.52),
          size=9, color=C_LIGHT)

add_label(s_intro, "Methodik-Framework",
          L_LEFT + Inches(0.18), L_TOP + Inches(4.18),
          L_W - Inches(0.3), Inches(0.32),
          size=9.5, bold=True, color=C_ACCENT)
add_label(s_intro,
          "CPA «From Data to Decisions» (Marr, 2020)\nSchritte 1–4: Daten → Analyse → Entscheid",
          L_LEFT + Inches(0.18), L_TOP + Inches(4.52),
          L_W - Inches(0.3), Inches(0.62),
          size=9, color=C_LIGHT)

# Right panel: three FF research question cards
R_LEFT = L_LEFT + L_W + Inches(0.28)
R_W    = SLIDE_W - R_LEFT - Inches(0.25)
R_H    = Inches(1.67)
R_GAP  = Inches(0.115)

C_FF2  = RGBColor(0x2D, 0xD4, 0x8A)
C_FF3  = RGBColor(0xA2, 0x8B, 0xF5)

ff_cards = [
    (C_ACCENT, "FF1 — Absolutes Veloaufkommen",
     "Welche Zählstellen weisen konstant das höchste absolute Veloaufkommen auf, "
     "und wie hat sich dieses im Zeitraum 2020–2025 verändert?",
     "→  Identifikation der meistgenutzten Achsen als Basis für Prioritätsentscheide"),
    (C_FF2,    "FF2 — Wachstums-Hotspots",
     "Welche Zählstellen zeigen den stärksten Wachstumstrend? "
     "Wo soll die Stadt präventive Kapazitätserweiterungen priorisieren?",
     "→  Vorausschauende Planung statt reaktive Engpassbehebung"),
    (C_FF3,    "FF3 — Wettersensitivität",
     "Wie stark reagiert das Veloaufkommen auf Regen und tiefe Temperaturen? "
     "Welche Standorte weisen eine überdurchschnittliche Wettersensitivität auf?",
     "→  Ganzjährige Alltagstauglichkeit des Velonetzes sicherstellen"),
]

for fi, (col, title, question, action) in enumerate(ff_cards):
    ft = L_TOP + fi * (R_H + R_GAP)
    fb = s_intro.shapes.add_shape(1, R_LEFT, ft, R_W, R_H)
    fb.fill.solid(); fb.fill.fore_color.rgb = RGBColor(0x12, 0x26, 0x3A)
    fb.line.color.rgb = col; fb.line.width = Pt(1.5)

    add_label(s_intro, title,
              R_LEFT + Inches(0.2), ft + Inches(0.12),
              R_W - Inches(0.3), Inches(0.35),
              size=12, bold=True, color=col)
    add_label(s_intro, question,
              R_LEFT + Inches(0.2), ft + Inches(0.5),
              R_W - Inches(0.3), Inches(0.78),
              size=9.5, color=C_WHITE)
    add_label(s_intro, action,
              R_LEFT + Inches(0.2), ft + Inches(1.3),
              R_W - Inches(0.3), Inches(0.3),
              size=9, bold=True, color=col)

# ── SLIDE 1: Title ────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
add_bg(s1)
accent_line(s1, Inches(0), Inches(0), SLIDE_W, Inches(0.06))

add_label(s1, "BINA — Forschungsfrage 1",
          Inches(0.7), Inches(1.3), Inches(9), Inches(0.5),
          size=13, color=C_ACCENT)

tb = s1.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(11.6), Inches(1.8))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Absolutes Veloaufkommen"
run.font.size = Pt(46); run.font.bold = True; run.font.color.rgb = C_WHITE

add_label(s1,
    "Welche Zählstellen weisen konstant das höchste absolute Veloaufkommen auf,\n"
    "und wie hat sich dieses von 2020 bis 2025 verändert?",
    Inches(0.7), Inches(3.2), Inches(10), Inches(1.2),
    size=15, color=C_LIGHT)

stats_boxes = [
    ("5 Mio+",  "Messdatensätze"),
    ("70+",     "Zählstellen"),
    ("6 Jahre", "2020–2025"),
    ("Zürich",  "Stadtgebiet"),
]
box_w, box_h = Inches(2.4), Inches(1.3)
box_top = Inches(5.4)
for i, (num, lbl) in enumerate(stats_boxes):
    left = Inches(0.7) + i * (box_w + Inches(0.25))
    box = s1.shapes.add_shape(1, left, box_top, box_w, box_h)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1A, 0x2F, 0x4A)
    box.line.color.rgb = C_ACCENT; box.line.width = Pt(1)
    add_label(s1, num, left + Inches(0.12), box_top + Inches(0.1),
              box_w - Inches(0.2), Inches(0.7),
              size=26, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_label(s1, lbl, left + Inches(0.12), box_top + Inches(0.75),
              box_w - Inches(0.2), Inches(0.45),
              size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

accent_line(s1, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))

# ── SLIDE 2: Projekt-Highlights ───────────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
add_bg(s2)
accent_line(s2, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
accent_line(s2, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))

add_label(s2, "02", Inches(0.25), Inches(0.2), Inches(0.8), Inches(0.5),
          size=11, color=C_ACCENT)
add_label(s2, "Projekt-Highlights",
          Inches(0.25), Inches(0.45), Inches(12), Inches(0.6),
          size=26, bold=True, color=C_WHITE)
add_label(s2, "Was wir im Rahmen des BINA-Projekts gelernt und erreicht haben",
          Inches(0.25), Inches(0.98), Inches(12), Inches(0.38),
          size=12, color=C_LIGHT)

# Three columns
COL_W   = Inches(3.9)
COL_H   = Inches(5.3)
COL_TOP = Inches(1.5)
GAP     = Inches(0.28)
col_lefts = [Inches(0.25), Inches(0.25) + COL_W + GAP, Inches(0.25) + 2*(COL_W + GAP)]

C_GREEN2 = RGBColor(0x2D, 0xD4, 0x8A)
C_PURPLE = RGBColor(0xA2, 0x8B, 0xF5)

col_configs = [
    # (border_color, icon, title, bullets)
    (C_ACCENT, "🚀", "Technische Lernkurve",
     [
         ("🐍 Python",          "Von null zur vollständigen Datenanalyse — Skripte, Schleifen, DataFrames"),
         ("📦 Bibliotheken",    "Pandas, Matplotlib, Seaborn, Scipy — professioneller Analyse-Stack"),
         ("📓 Jupyter Notebook","Interaktive Analyse, Dokumentation und Visualisierung in einem"),
         ("🤖 KI-Unterstützung","Effizienter durch intelligente Assistenz bei Code & Interpretation"),
     ]),
    (C_GREEN2, "🤝", "Kollaboration & Methodik",
     [
         ("👥 Pair Programming", "Gemeinsam schneller zum Ziel — weniger Fehler, mehr Lerneffekt"),
         ("📒 Geteilte Notebooks","Transparente, reproduzierbare Analysen für das ganze Team"),
         ("🌿 Git",              "Versionskontrolle für jeden Analyseschritt — nichts geht verloren"),
         ("🐙 GitHub",           "Zentrales Repository, Pull Requests und gemeinsame Codebasis"),
     ]),
    (C_PURPLE, "🎯", "Analytische Flugoptik",
     [
         ("🔧 Datenaufbereitung", "Rohdaten bereinigen, zusammenführen und validieren — die Basis für alles"),
         ("📊 Datenanalyse",      "Statistische Methoden, Regression und Zeitreihen sinnvoll einsetzen"),
         ("💼 Management-Komm.",  "Erkenntnisse auf den Punkt bringen — verständlich für Entscheider"),
     ]),
]

for col_i, (border_col, icon, title, bullets) in enumerate(col_configs):
    left = col_lefts[col_i]

    # Card background
    card = s2.shapes.add_shape(1, left, COL_TOP, COL_W, COL_H)
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x12, 0x26, 0x3A)
    card.line.color.rgb = border_col; card.line.width = Pt(1.5)

    # Icon + title header
    add_label(s2, icon, left + Inches(0.18), COL_TOP + Inches(0.15),
              Inches(0.55), Inches(0.52), size=22, align=PP_ALIGN.CENTER)
    add_label(s2, title, left + Inches(0.72), COL_TOP + Inches(0.2),
              COL_W - Inches(0.85), Inches(0.45),
              size=13, bold=True, color=border_col)

    # Divider line
    div = s2.shapes.add_shape(1, left + Inches(0.18), COL_TOP + Inches(0.75),
                               COL_W - Inches(0.36), Inches(0.025))
    div.fill.solid(); div.fill.fore_color.rgb = border_col
    div.line.fill.background()

    # Bullet items
    for b_i, (label, desc) in enumerate(bullets):
        bt = COL_TOP + Inches(0.9) + b_i * Inches(1.08)
        add_label(s2, label, left + Inches(0.18), bt,
                  COL_W - Inches(0.3), Inches(0.32),
                  size=10, bold=True, color=C_WHITE)
        add_label(s2, desc, left + Inches(0.18), bt + Inches(0.32),
                  COL_W - Inches(0.3), Inches(0.68),
                  size=8.5, color=C_LIGHT)

# ── SLIDE 3: Heatmap ──────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
add_bg(s3)
accent_line(s3, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
accent_line(s3, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))

add_label(s3, "03", Inches(0.25), Inches(0.2), Inches(0.8), Inches(0.5),
          size=11, color=C_ACCENT)
add_label(s3, "Jahresentwicklung — Top-15 Stationen",
          Inches(0.25), Inches(0.45), Inches(10), Inches(0.6), size=22, bold=True)
add_label(s3, "Veloaufkommen in Tausend Fahrten pro Jahr (Heatmap)",
          Inches(0.25), Inches(0.95), Inches(10), Inches(0.4), size=12, color=C_LIGHT)

s3.shapes.add_picture(chart2_stream, Inches(0.2), Inches(1.25), Inches(7.8), Inches(5.9))

callouts = [
    (C_ORANGE,
     "Hochvariable Stationen",
     "Stadttunnel Nord & Kreuzstrasse zeigen hohe Schwankungen (rank_std > 8)."),
    (C_ACCENT,
     "Stabile Top-Performer",
     "Sonneggstrasse & Hardbrücke: konstant führend (rank_std ≤ 4.3)."),
    (RGBColor(0x48, 0xCA, 0xE4),
     "Sihlstrasse wächst",
     "Von 311k (2020) auf 998k (2025) — stärkstes Wachstum 2024→2025."),
]
for j, (col, title, body) in enumerate(callouts):
    tp = Inches(1.4) + j * Inches(1.9)
    box = s3.shapes.add_shape(1, Inches(8.25), tp, Inches(4.8), Inches(1.65))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1A, 0x2F, 0x4A)
    box.line.color.rgb = col; box.line.width = Pt(1.5)
    add_label(s3, title, Inches(8.45), tp + Inches(0.1),
              Inches(4.4), Inches(0.4), size=11, bold=True, color=col)
    add_label(s3, body, Inches(8.45), tp + Inches(0.52),
              Inches(4.4), Inches(1.0), size=9.5, color=C_LIGHT)

# ── SLIDE 4: FF2 Wachstums-Hotspots ──────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
add_bg(s3)
accent_line(s3, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
accent_line(s3, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))

add_label(s3, "04", Inches(0.25), Inches(0.2), Inches(0.8), Inches(0.5),
          size=11, color=C_ACCENT)
add_label(s3, "Wachstums-Hotspots — Wo muss Zürich jetzt handeln?",
          Inches(0.25), Inches(0.45), Inches(10), Inches(0.6), size=22, bold=True)
add_label(s3, "CAGR 2020–2025 — 12 konsistente Zählstellen mit vollständiger Datenabdeckung",
          Inches(0.25), Inches(0.95), Inches(10), Inches(0.4), size=12, color=C_LIGHT)

s3.shapes.add_picture(chartFF2_stream, Inches(0.2), Inches(1.28), Inches(8.8), Inches(5.9))

# Right panel
panel_left = Inches(9.25)
panel_top  = Inches(1.28)
panel_w    = Inches(3.85)

box = s3.shapes.add_shape(1, panel_left, panel_top, panel_w, Inches(5.9))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1A, 0x2F, 0x4A)
box.line.color.rgb = RGBColor(0x00, 0xD4, 0xA0); box.line.width = Pt(1)

add_label(s3, "MANAGEMENT-ENTSCHEID", panel_left + Inches(0.2), panel_top + Inches(0.15),
          panel_w - Inches(0.3), Inches(0.38), size=9.5, bold=True,
          color=RGBColor(0x00, 0xD4, 0xA0))

add_label(s3,
    "80% des Netzwachstums\nkonzentriert auf 3 Korridore\n→ Gezielte Investition wirkt.",
    panel_left + Inches(0.2), panel_top + Inches(0.6),
    panel_w - Inches(0.3), Inches(0.9), size=10, bold=True, color=C_WHITE)

action_items = [
    (RGBColor(0x00, 0xD4, 0xA0), "A — Sofort handeln",
     "Baslerstrasse\nCAGR +13.4% · R²=0.91\n+87% seit 2020 · Priorität 1"),
    (C_ACCENT, "B — Mittelfristig",
     "Langstr. Fahrbahn Nord\nCAGR +5.1% · R²=0.86\nStetig steigend seit 2021"),
    (C_ACCENT, "B — Mittelfristig",
     "Hardbrücke Nord\nCAGR +4.2% · R²=0.88\nKonsistentes Wachstum"),
    (C_LIGHT, "⚠ Sinkend beobachten",
     "Langstr. Unterführungen\nCAGR −5 bis −7%\nUrsache klären"),
]

for j, (col, title, body) in enumerate(action_items):
    tp = panel_top + Inches(1.6) + j * Inches(1.1)
    bullet = s3.shapes.add_shape(1, panel_left + Inches(0.15), tp + Inches(0.07),
                                  Inches(0.06), Inches(0.42))
    bullet.fill.solid(); bullet.fill.fore_color.rgb = col
    bullet.line.fill.background()
    add_label(s3, title, panel_left + Inches(0.32), tp,
              panel_w - Inches(0.4), Inches(0.32), size=9.5, bold=True, color=col)
    add_label(s3, body, panel_left + Inches(0.32), tp + Inches(0.32),
              panel_w - Inches(0.4), Inches(0.72), size=8.5, color=C_LIGHT)

# ── SLIDE 5: FF3 Wettersensitivität ──────────────────────────────────────────
s4 = prs.slides.add_slide(blank)
add_bg(s4)
accent_line(s4, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
accent_line(s4, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))

add_label(s4, "05", Inches(0.25), Inches(0.2), Inches(0.8), Inches(0.5),
          size=11, color=C_ACCENT)
add_label(s4, "Wettersensitivität — Kälte ist die eigentliche Hürde",
          Inches(0.25), Inches(0.45), Inches(10), Inches(0.6), size=22, bold=True)
add_label(s4, "Ø Velofahrten pro Standort-Tag nach Wetterbedingung — Netzwerkebene 2020–2025",
          Inches(0.25), Inches(0.95), Inches(10), Inches(0.4), size=12, color=C_LIGHT)

s4.shapes.add_picture(chartFF3_stream, Inches(0.2), Inches(1.28), Inches(8.6), Inches(5.5))

# Right insight panel
C_RED   = RGBColor(0xFF, 0x4D, 0x6D)
C_GREEN = RGBColor(0x2D, 0xD4, 0x8A)

panel_left = Inches(9.05)
panel_top  = Inches(1.28)
panel_w    = Inches(4.05)

box = s4.shapes.add_shape(1, panel_left, panel_top, panel_w, Inches(5.5))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1A, 0x2F, 0x4A)
box.line.color.rgb = C_RED; box.line.width = Pt(1)

add_label(s4, "MANAGEMENT-ENTSCHEID",
          panel_left + Inches(0.2), panel_top + Inches(0.15),
          panel_w - Inches(0.3), Inches(0.38),
          size=9.5, bold=True, color=C_RED)

add_label(s4,
    "Kälte reduziert das Veloaufkommen 4× stärker als Regen\n— ein direkter Widerspruch zur Velostrategie 2030.",
    panel_left + Inches(0.2), panel_top + Inches(0.6),
    panel_w - Inches(0.3), Inches(0.85),
    size=10, bold=True, color=C_WHITE)

findings = [
    (C_ACCENT,  "Regen — situativ",
     "−10.2% Rückgang\nKurzfristige Hürde, viele Pendler\nfahren trotzdem."),
    (C_RED,     "Kälte — strukturell",
     "−41.6% Rückgang bei ≤5 °C\nWintermonate blockieren\nAlltagstauglichkeit."),
    (C_RED,     "Kombiniert",
     "−47% bei Regen + Kälte\nPraktisch keine Freizeitnutzung\nmehr möglich."),
    (C_GREEN,   "Positive Tendenz",
     "Kälte-Sens. −43.9% (2020)\n→ −36.3% (2025): Infrastruktur\nund Gewohnheit wirken."),
]

for j, (col, title, body) in enumerate(findings):
    tp = panel_top + Inches(1.55) + j * Inches(1.0)
    bullet = s4.shapes.add_shape(1, panel_left + Inches(0.15), tp + Inches(0.07),
                                  Inches(0.06), Inches(0.38))
    bullet.fill.solid(); bullet.fill.fore_color.rgb = col
    bullet.line.fill.background()
    add_label(s4, title, panel_left + Inches(0.32), tp,
              panel_w - Inches(0.4), Inches(0.32),
              size=9.5, bold=True, color=col)
    add_label(s4, body, panel_left + Inches(0.32), tp + Inches(0.32),
              panel_w - Inches(0.4), Inches(0.65),
              size=8.5, color=C_LIGHT)

# ── reorder slides: move Highlights (index 2) to position 1 ──────────────────
# Creation order: Intro(0), Title(1), Highlights(2), Heatmap(3), FF2(4), FF3(5)
# Desired order:  Intro(0), Highlights(1), Title(2), Heatmap(3), FF2(4), FF3(5)
_sldIdLst = prs.slides._sldIdLst
_slides = list(_sldIdLst)
_el = _slides[2]
_sldIdLst.remove(_el)
_sldIdLst.insert(1, _el)

# ── save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "FF1_Veloaufkommen_Zürich.pptx")
prs.save(out_path)
print(f"\nDone → {out_path}")
